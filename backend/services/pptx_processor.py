from pptx import Presentation

def extract_pptx_text(file_path):
    prs = Presentation(file_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)
